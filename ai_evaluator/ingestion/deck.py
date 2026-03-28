import os
import fitz  # PyMuPDF
from pptx import Presentation


def extract_pdf(file_path):
    """
    Extract text from PDF file (each page treated as a slide)
    """
    slides = []

    try:
        doc = fitz.open(file_path)

        for i, page in enumerate(doc):
            text = page.get_text().strip()

            slides.append({
                "source": f"slide_{i}",
                "type": "deck",
                "format": "pdf",
                "content": text if text else "",
            })

    except Exception as e:
        raise RuntimeError(f"Error processing PDF: {str(e)}")

    return slides


def extract_pptx(file_path):
    """
    Extract text from PPT/PPTX file
    """
    slides = []

    try:
        prs = Presentation(file_path)

        for i, slide in enumerate(prs.slides):
            text_runs = []
            title = ""

            # Extract title if present
            try:
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
            except Exception:
                title = ""

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        text_runs.append(text)

            full_text = "\n".join(text_runs)

            slides.append({
                "source": f"slide_{i}",
                "type": "deck",
                "format": "pptx",
                "title": title,
                "content": f"Title: {title}\n{full_text}" if title else full_text,
            })

    except Exception as e:
        raise RuntimeError(f"Error processing PPT/PPTX: {str(e)}")

    return slides


def extract_deck_text(file_path):
    """
    Main function to extract deck content from PDF or PPT/PPTX

    Returns:
        List[Dict] with keys:
        - source (slide index)
        - type (deck)
        - format (pdf/pptx)
        - content (text)
    """

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf(file_path)

    elif ext in [".pptx", ".ppt"]:
        return extract_pptx(file_path)

    else:
        raise ValueError(
            f"Unsupported file format: {ext}. Supported formats: .pdf, .ppt, .pptx"
        )